"""
PPTX Generator - Gerador de Apresentações PowerPoint

Este módulo é responsável por gerar apresentações PPTX automaticamente:
- Carrega templates PPTX com placeholders
- Substitui placeholders por dados reais
- Formata slides conforme necessário
- Salva apresentação final

Instruções para personalização:
1. Crie templates PPTX com placeholders no formato {{nome_placeholder}}
2. Adicione novos tipos de formatação conforme necessário
3. Customize os estilos de slide conforme sua identidade visual
4. Adicione novos tipos de elementos (gráficos, tabelas, etc.)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import pandas as pd
import re
import os
from typing import Dict, List, Optional, Union, Any, Tuple
from pathlib import Path
from loguru import logger
import json
from datetime import datetime
import shutil


class PPTXGenerator:
    """
    Classe principal para geração de apresentações PPTX.
    
    Funcionalidades:
    - Carregamento de templates PPTX
    - Substituição de placeholders por dados reais
    - Formatação de slides
    - Geração de gráficos e tabelas
    - Salvamento da apresentação final
    """
    
    def __init__(self, config: Optional[Dict] = None):
        """
        Inicializa o PPTXGenerator com configurações opcionais.
        
        Args:
            config (Dict, optional): Configurações personalizadas
        """
        self.config = config or {}
        
        # Configurações padrão - EDITE AQUI conforme necessário
        self.default_config = {
            'placeholder_pattern': r'{{(.*?)}}',  # Padrão para identificar placeholders
            'default_font': 'Calibri',  # Fonte padrão
            'default_font_size': Pt(12),  # Tamanho de fonte padrão
            'title_font_size': Pt(24),  # Tamanho de fonte para títulos
            'subtitle_font_size': Pt(18),  # Tamanho de fonte para subtítulos
            'chart_width': Inches(6),  # Largura padrão para gráficos
            'chart_height': Inches(4),  # Altura padrão para gráficos
            'table_style': 'LightGrid',  # Estilo padrão para tabelas
            'image_width': Inches(4),  # Largura padrão para imagens
            'image_height': Inches(3),  # Altura padrão para imagens
            'default_colors': [  # Cores padrão para gráficos
                RGBColor(31, 119, 180),  # Azul
                RGBColor(255, 127, 14),  # Laranja
                RGBColor(44, 160, 44),   # Verde
                RGBColor(214, 39, 40),   # Vermelho
                RGBColor(148, 103, 189), # Roxo
                RGBColor(140, 86, 75),   # Marrom
                RGBColor(227, 119, 194), # Rosa
                RGBColor(127, 127, 127), # Cinza
                RGBColor(188, 189, 34),  # Verde-limão
                RGBColor(23, 190, 207)   # Ciano
            ]
        }
        
        # Mescla configurações
        self.settings = {**self.default_config, **self.config}
        
        # Inicializa variáveis
        self.presentation = None
        self.template_path = None
        self.placeholder_pattern = re.compile(self.settings['placeholder_pattern'])
        
        logger.info("PPTXGenerator inicializado com sucesso")
    
    def load_template(self, template_path: str) -> bool:
        """
        Carrega um template PPTX.
        
        Args:
            template_path (str): Caminho para o arquivo de template PPTX
            
        Returns:
            bool: True se o template foi carregado com sucesso
            
        Exemplo de uso:
            generator.load_template('templates/relatorio.pptx')
        """
        try:
            logger.info(f"Carregando template PPTX: {template_path}")
            
            if not os.path.exists(template_path):
                raise FileNotFoundError(f"Template não encontrado: {template_path}")
            
            self.presentation = Presentation(template_path)
            self.template_path = template_path
            
            # Conta placeholders no template
            placeholder_count = self._count_placeholders()
            
            logger.info(f"Template carregado com sucesso: {len(self.presentation.slides)} slides, {placeholder_count} placeholders")
            return True
            
        except Exception as e:
            logger.error(f"Erro ao carregar template: {str(e)}")
            return False
    
    def create_new_presentation(self) -> bool:
        """
        Cria uma nova apresentação em branco.
        
        Returns:
            bool: True se a apresentação foi criada com sucesso
            
        Exemplo de uso:
            generator.create_new_presentation()
        """
        try:
            logger.info("Criando nova apresentação PPTX")
            
            self.presentation = Presentation()
            self.template_path = None
            
            logger.info("Nova apresentação criada com sucesso")
            return True
            
        except Exception as e:
            logger.error(f"Erro ao criar nova apresentação: {str(e)}")
            return False
    
    def replace_placeholders(self, data: Dict) -> int:
        """
        Substitui placeholders nos slides por dados reais.
        
        Args:
            data (Dict): Dicionário com dados para substituir os placeholders
            
        Returns:
            int: Número de placeholders substituídos
            
        Exemplo de uso:
            data = {'total': 1500, 'media_2024': 125.5, 'media_2025': 142.8}
            generator.replace_placeholders(data)
        """
        try:
            logger.info("Substituindo placeholders")
            
            if not self.presentation:
                raise ValueError("Nenhuma apresentação carregada")
            
            replaced_count = 0
            
            # Percorre todos os slides
            for slide in self.presentation.slides:
                # Percorre todas as formas no slide
                for shape in slide.shapes:
                    # Verifica se a forma tem texto
                    if hasattr(shape, "text_frame"):
                        # Substitui placeholders no texto
                        replaced_count += self._replace_in_text_frame(shape.text_frame, data)
                    
                    # Verifica se é uma tabela
                    if shape.has_table:
                        # Substitui placeholders na tabela
                        replaced_count += self._replace_in_table(shape.table, data)
            
            logger.info(f"Substituição concluída: {replaced_count} placeholders substituídos")
            return replaced_count
            
        except Exception as e:
            logger.error(f"Erro ao substituir placeholders: {str(e)}")
            return 0
    
    def add_slide(self, layout_name: Optional[str] = None) -> int:
        """
        Adiciona um novo slide à apresentação.
        
        Args:
            layout_name (str, optional): Nome do layout a ser usado
            
        Returns:
            int: Índice do slide adicionado
            
        Exemplo de uso:
            slide_idx = generator.add_slide('Title and Content')
        """
        try:
            logger.info(f"Adicionando novo slide com layout: {layout_name or 'padrão'}")
            
            if not self.presentation:
                raise ValueError("Nenhuma apresentação carregada")
            
            # Se não especificou layout, usa o primeiro disponível
            if layout_name is None:
                slide_layout = self.presentation.slide_layouts[0]
            else:
                # Procura layout pelo nome
                slide_layout = None
                for layout in self.presentation.slide_layouts:
                    if layout.name == layout_name:
                        slide_layout = layout
                        break
                
                # Se não encontrou, usa o primeiro
                if slide_layout is None:
                    logger.warning(f"Layout '{layout_name}' não encontrado, usando padrão")
                    slide_layout = self.presentation.slide_layouts[0]
            
            # Adiciona slide
            slide = self.presentation.slides.add_slide(slide_layout)
            slide_idx = len(self.presentation.slides) - 1
            
            logger.info(f"Slide adicionado com sucesso: índice {slide_idx}")
            return slide_idx
            
        except Exception as e:
            logger.error(f"Erro ao adicionar slide: {str(e)}")
            return -1
    
    def add_chart(self, slide_idx: int, chart_data: Dict, chart_type: str = 'column', 
                 position: Tuple[float, float] = (2, 2), size: Tuple[float, float] = (6, 4)) -> bool:
        """
        Adiciona um gráfico a um slide.
        
        Args:
            slide_idx (int): Índice do slide
            chart_data (Dict): Dados para o gráfico
            chart_type (str): Tipo de gráfico (column, bar, line, pie)
            position (Tuple[float, float]): Posição (x, y) em polegadas
            size (Tuple[float, float]): Tamanho (largura, altura) em polegadas
            
        Returns:
            bool: True se o gráfico foi adicionado com sucesso
            
        Exemplo de uso:
            data = {
                'categories': ['Jan', 'Fev', 'Mar'],
                'series': [
                    {'name': '2024', 'values': [10, 20, 30]},
                    {'name': '2025', 'values': [15, 25, 35]}
                ]
            }
            generator.add_chart(0, data, 'column', (2, 2), (6, 4))
        """
        try:
            logger.info(f"Adicionando gráfico ao slide {slide_idx}")
            
            if not self.presentation:
                raise ValueError("Nenhuma apresentação carregada")
            
            if slide_idx < 0 or slide_idx >= len(self.presentation.slides):
                raise ValueError(f"Índice de slide inválido: {slide_idx}")
            
            # Obtém slide
            slide = self.presentation.slides[slide_idx]
            
            # Cria dados do gráfico
            chart_data_obj = CategoryChartData()
            
            # Adiciona categorias
            categories = chart_data.get('categories', [])
            chart_data_obj.categories = categories
            
            # Adiciona séries
            for series in chart_data.get('series', []):
                series_name = series.get('name', '')
                series_values = series.get('values', [])
                chart_data_obj.add_series(series_name, series_values)
            
            # Mapeia tipo de gráfico
            chart_type_map = {
                'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'bar': XL_CHART_TYPE.BAR_CLUSTERED,
                'line': XL_CHART_TYPE.LINE,
                'pie': XL_CHART_TYPE.PIE,
                'area': XL_CHART_TYPE.AREA,
                'scatter': XL_CHART_TYPE.XY_SCATTER,
                'doughnut': XL_CHART_TYPE.DOUGHNUT
            }
            
            xl_chart_type = chart_type_map.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
            
            # Adiciona gráfico ao slide
            x, y = position
            width, height = size
            chart = slide.shapes.add_chart(
                xl_chart_type,
                Inches(x),
                Inches(y),
                Inches(width),
                Inches(height),
                chart_data_obj
            ).chart
            
            # Formata gráfico
            chart.has_legend = True
            chart.legend.position = 1  # Posição à direita
            chart.legend.include_in_layout = False
            
            logger.info(f"Gráfico adicionado com sucesso ao slide {slide_idx}")
            return True
            
        except Exception as e:
            logger.error(f"Erro ao adicionar gráfico: {str(e)}")
            return False
    
    def add_table(self, slide_idx: int, table_data: List[List], 
                 position: Tuple[float, float] = (2, 2), size: Tuple[float, float] = (6, 3)) -> bool:
        """
        Adiciona uma tabela a um slide.
        
        Args:
            slide_idx (int): Índice do slide
            table_data (List[List]): Dados para a tabela (lista de linhas, cada linha é uma lista de células)
            position (Tuple[float, float]): Posição (x, y) em polegadas
            size (Tuple[float, float]): Tamanho (largura, altura) em polegadas
            
        Returns:
            bool: True se a tabela foi adicionada com sucesso
            
        Exemplo de uso:
            data = [
                ['Métrica', '2024', '2025', 'Variação'],
                ['Total', '1500', '1800', '+20%'],
                ['Média', '125', '150', '+20%']
            ]
            generator.add_table(0, data, (2, 2), (6, 2))
        """
        try:
            logger.info(f"Adicionando tabela ao slide {slide_idx}")
            
            if not self.presentation:
                raise ValueError("Nenhuma apresentação carregada")
            
            if slide_idx < 0 or slide_idx >= len(self.presentation.slides):
                raise ValueError(f"Índice de slide inválido: {slide_idx}")
            
            if not table_data or not table_data[0]:
                raise ValueError("Dados de tabela vazios")
            
            # Obtém slide
            slide = self.presentation.slides[slide_idx]
            
            # Determina dimensões da tabela
            rows = len(table_data)
            cols = len(table_data[0])
            
            # Adiciona tabela ao slide
            x, y = position
            width, height = size
            shape = slide.shapes.add_table(
                rows,
                cols,
                Inches(x),
                Inches(y),
                Inches(width),
                Inches(height)
            )
            
            table = shape.table
            
            # Preenche tabela com dados
            for i, row_data in enumerate(table_data):
                for j, cell_text in enumerate(row_data):
                    if j < cols:  # Garante que não ultrapasse o número de colunas
                        cell = table.cell(i, j)
                        cell.text = str(cell_text)
                        
                        # Formata cabeçalho
                        if i == 0:
                            cell.text_frame.paragraphs[0].font.bold = True
                            cell.text_frame.paragraphs[0].font.size = Pt(14)
                            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            logger.info(f"Tabela adicionada com sucesso ao slide {slide_idx}")
            return True
            
        except Exception as e:
            logger.error(f"Erro ao adicionar tabela: {str(e)}")
            return False
    
    def add_image(self, slide_idx: int, image_path: str, 
                 position: Tuple[float, float] = (2, 2), size: Optional[Tuple[float, float]] = None) -> bool:
        """
        Adiciona uma imagem a um slide.
        
        Args:
            slide_idx (int): Índice do slide
            image_path (str): Caminho para o arquivo de imagem
            position (Tuple[float, float]): Posição (x, y) em polegadas
            size (Tuple[float, float], optional): Tamanho (largura, altura) em polegadas
            
        Returns:
            bool: True se a imagem foi adicionada com sucesso
            
        Exemplo de uso:
            generator.add_image(0, 'images/grafico.png', (2, 2), (4, 3))
        """
        try:
            logger.info(f"Adicionando imagem ao slide {slide_idx}")
            
            if not self.presentation:
                raise ValueError("Nenhuma apresentação carregada")
            
            if slide_idx < 0 or slide_idx >= len(self.presentation.slides):
                raise ValueError(f"Índice de slide inválido: {slide_idx}")
            
            if not os.path.exists(image_path):
                raise FileNotFoundError(f"Imagem não encontrada: {image_path}")
            
            # Obtém slide
            slide = self.presentation.slides[slide_idx]
            
            # Adiciona imagem ao slide
            x, y = position
            
            if size:
                width, height = size
                slide.shapes.add_picture(
                    image_path,
                    Inches(x),
                    Inches(y),
                    width=Inches(width),
                    height=Inches(height)
                )
            else:
                slide.shapes.add_picture(
                    image_path,
                    Inches(x),
                    Inches(y)
                )
            
            logger.info(f"Imagem adicionada com sucesso ao slide {slide_idx}")
            return True
            
        except Exception as e:
            logger.error(f"Erro ao adicionar imagem: {str(e)}")
            return False
    
    def save_presentation(self, output_path: str) -> bool:
        """
        Salva a apresentação em um arquivo.
        
        Args:
            output_path (str): Caminho para salvar o arquivo PPTX
            
        Returns:
            bool: True se a apresentação foi salva com sucesso
            
        Exemplo de uso:
            generator.save_presentation('output/relatorio_final.pptx')
        """
        try:
            logger.info(f"Salvando apresentação: {output_path}")
            
            if not self.presentation:
                raise ValueError("Nenhuma apresentação carregada")
            
            # Cria diretório se não existir
            os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
            
            # Salva apresentação
            self.presentation.save(output_path)
            
            logger.info(f"Apresentação salva com sucesso: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"Erro ao salvar apresentação: {str(e)}")
            return False
    
    def get_available_layouts(self) -> List[str]:
        """
        Retorna os layouts disponíveis na apresentação.
        
        Returns:
            List[str]: Lista de nomes de layouts disponíveis
            
        Exemplo de uso:
            layouts = generator.get_available_layouts()
        """
        try:
            if not self.presentation:
                raise ValueError("Nenhuma apresentação carregada")
            
            layouts = [layout.name for layout in self.presentation.slide_layouts]
            return layouts
            
        except Exception as e:
            logger.error(f"Erro ao obter layouts: {str(e)}")
            return []
    
    def get_placeholders_list(self) -> List[str]:
        """
        Retorna a lista de placeholders encontrados no template.
        
        Returns:
            List[str]: Lista de placeholders encontrados
            
        Exemplo de uso:
            placeholders = generator.get_placeholders_list()
        """
        try:
            if not self.presentation:
                raise ValueError("Nenhuma apresentação carregada")
            
            placeholders = set()
            
            # Percorre todos os slides
            for slide in self.presentation.slides:
                # Percorre todas as formas no slide
                for shape in slide.shapes:
                    # Verifica se a forma tem texto
                    if hasattr(shape, "text_frame"):
                        # Extrai placeholders do texto
                        self._extract_placeholders(shape.text_frame.text, placeholders)
                    
                    # Verifica se é uma tabela
                    if shape.has_table:
                        # Extrai placeholders da tabela
                        for row in shape.table.rows:
                            for cell in row.cells:
                                self._extract_placeholders(cell.text, placeholders)
            
            return sorted(list(placeholders))
            
        except Exception as e:
            logger.error(f"Erro ao obter lista de placeholders: {str(e)}")
            return []
    
    def create_template_from_data(self, data: Dict, output_path: str) -> bool:
        """
        Cria um template PPTX a partir de dados de exemplo.
        
        Args:
            data (Dict): Dados de exemplo para criar o template
            output_path (str): Caminho para salvar o template
            
        Returns:
            bool: True se o template foi criado com sucesso
            
        Exemplo de uso:
            data = {
                'titulo': 'Relatório de Vendas',
                'subtitulo': 'Comparativo 2024-2025',
                'total_2024': 1500,
                'total_2025': 1800,
                'crescimento': '20%'
            }
            generator.create_template_from_data(data, 'templates/template_relatorio.pptx')
        """
        try:
            logger.info(f"Criando template a partir de dados: {output_path}")
            
            # Cria nova apresentação
            self.create_new_presentation()
            
            # Adiciona slide de título
            slide_idx = self.add_slide()
            slide = self.presentation.slides[slide_idx]
            
            # Adiciona título e subtítulo
            if hasattr(slide, "shapes"):
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        if shape.text_frame.text == "Click to add title":
                            shape.text_frame.text = "{{titulo}}"
                        elif shape.text_frame.text == "Click to add subtitle":
                            shape.text_frame.text = "{{subtitulo}}"
            
            # Adiciona slide de conteúdo
            slide_idx = self.add_slide()
            slide = self.presentation.slides[slide_idx]
            
            # Adiciona título
            if hasattr(slide, "shapes"):
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        if shape.text_frame.text == "Click to add title":
                            shape.text_frame.text = "Resultados"
            
            # Adiciona tabela com placeholders
            table_data = [
                ['Métrica', '2024', '2025', 'Variação'],
                ['Total', '{{total_2024}}', '{{total_2025}}', '{{crescimento}}']
            ]
            self.add_table(slide_idx, table_data, (2, 2), (6, 1.5))
            
            # Salva template
            self.save_presentation(output_path)
            
            logger.info(f"Template criado com sucesso: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"Erro ao criar template: {str(e)}")
            return False
    
    # Métodos auxiliares privados
    
    def _count_placeholders(self) -> int:
        """Conta o número de placeholders no template."""
        if not self.presentation:
            return 0
        
        count = 0
        placeholders = set()
        
        # Percorre todos os slides
        for slide in self.presentation.slides:
            # Percorre todas as formas no slide
            for shape in slide.shapes:
                # Verifica se a forma tem texto
                if hasattr(shape, "text_frame"):
                    # Extrai placeholders do texto
                    self._extract_placeholders(shape.text_frame.text, placeholders)
                
                # Verifica se é uma tabela
                if shape.has_table:
                    # Extrai placeholders da tabela
                    for row in shape.table.rows:
                        for cell in row.cells:
                            self._extract_placeholders(cell.text, placeholders)
        
        return len(placeholders)
    
    def _extract_placeholders(self, text: str, placeholders: set) -> None:
        """Extrai placeholders de um texto e adiciona ao conjunto."""
        if not text:
            return
        
        matches = self.placeholder_pattern.findall(text)
        for match in matches:
            placeholders.add(match)
    
    def _replace_in_text_frame(self, text_frame, data: Dict) -> int:
        """Substitui placeholders em um frame de texto."""
        replaced_count = 0
        
        if not text_frame.text:
            return replaced_count
        
        # Substitui no texto principal
        original_text = text_frame.text
        new_text = original_text
        
        for match in self.placeholder_pattern.findall(original_text):
            if match in data:
                placeholder = f"{{{{{match}}}}}"
                new_text = new_text.replace(placeholder, str(data[match]))
                replaced_count += 1
        
        if new_text != original_text:
            text_frame.text = new_text
        
        # Substitui em cada parágrafo
        for paragraph in text_frame.paragraphs:
            original_text = paragraph.text
            new_text = original_text
            
            for match in self.placeholder_pattern.findall(original_text):
                if match in data:
                    placeholder = f"{{{{{match}}}}}"
                    new_text = new_text.replace(placeholder, str(data[match]))
                    replaced_count += 1
            
            if new_text != original_text:
                paragraph.text = new_text
        
        return replaced_count
    
    def _replace_in_table(self, table, data: Dict) -> int:
        """Substitui placeholders em uma tabela."""
        replaced_count = 0
        
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    original_text = paragraph.text
                    new_text = original_text
                    
                    for match in self.placeholder_pattern.findall(original_text):
                        if match in data:
                            placeholder = f"{{{{{match}}}}}"
                            new_text = new_text.replace(placeholder, str(data[match]))
                            replaced_count += 1
                    
                    if new_text != original_text:
                        paragraph.text = new_text
        
        return replaced_count


# Exemplo de uso e testes
if __name__ == "__main__":
    # Inicializa o PPTXGenerator
    generator = PPTXGenerator()
    
    # Exemplo de criação de template
    try:
        # Cria template de exemplo
        data = {
            'titulo': 'Relatório de Analytics',
            'subtitulo': 'Comparativo 2024-2025',
            'total_2024': 1500,
            'total_2025': 1800,
            'crescimento': '20%'
        }
        
        # generator.create_template_from_data(data, 'templates/template_exemplo.pptx')
        # print("Template criado com sucesso!")
        pass
    except Exception as e:
        print(f"Erro ao criar template: {e}")
    
    print("PPTXGenerator implementado com sucesso!")

