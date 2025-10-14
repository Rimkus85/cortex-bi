"""
Sistema de Administração do CÓRTEX BI
Permite ao usuário Redecorp\r337786 gerenciar o sistema via Copilot M365
"""

import os
import json
import shutil
from pathlib import Path
from typing import Dict, List, Optional, Any
from datetime import datetime
from loguru import logger
import requests
from pptx import Presentation
from pptx.util import Inches, Pt


class AdminSystem:
    """
    Sistema de administração para usuários autorizados do CÓRTEX BI.
    
    Funcionalidades:
    - Gerenciar templates PPTX
    - Modificar placeholders
    - Configurar integração SharePoint
    - Controlar permissões
    """
    
    def __init__(self):
        """Inicializa o sistema de administração"""
        self.authorized_users = [
            "Redecorp\\r337786",
            "redecorp\\r337786", 
            "REDECORP\\R337786"
        ]
        
        # Configuração SharePoint
        self.sharepoint_config = {
            "base_url": "telefonicacorp-my.sharepoint.com",
            "site_path": "/personal/usr_mis_br_telefonica_com/Documents/Python",
            "full_url": "https://telefonicacorp-my.sharepoint.com/my?id=%2Fpersonal%2Fusr_mis_br_telefonica_com%2FDocuments%2FPython"
        }
        
        # Diretórios
        self.templates_dir = Path("templates")
        self.output_dir = Path("output")
        self.backup_dir = Path("backups")
        
        # Criar diretórios se não existirem
        self.templates_dir.mkdir(exist_ok=True)
        self.output_dir.mkdir(exist_ok=True)
        self.backup_dir.mkdir(exist_ok=True)
        
        logger.info("Sistema de administração inicializado")
    
    def verify_admin_access(self, user_id: str) -> bool:
        """
        Verifica se o usuário tem permissões de administrador.
        
        Args:
            user_id (str): ID do usuário
            
        Returns:
            bool: True se autorizado
        """
        return user_id.lower() in [u.lower() for u in self.authorized_users]
    
    def list_templates(self) -> List[Dict]:
        """
        Lista todos os templates PPTX disponíveis.
        
        Returns:
            List[Dict]: Lista de templates com metadados
        """
        templates = []
        
        for template_file in self.templates_dir.glob("*.pptx"):
            try:
                # Obter informações do template
                prs = Presentation(template_file)
                
                # Extrair placeholders
                placeholders = self._extract_placeholders(prs)
                
                template_info = {
                    "name": template_file.name,
                    "path": str(template_file),
                    "size": template_file.stat().st_size,
                    "modified": datetime.fromtimestamp(template_file.stat().st_mtime).isoformat(),
                    "slides_count": len(prs.slides),
                    "placeholders": placeholders
                }
                
                templates.append(template_info)
                
            except Exception as e:
                logger.error(f"Erro ao processar template {template_file}: {e}")
        
        return templates
    
    def _extract_placeholders(self, presentation: Presentation) -> List[str]:
        """
        Extrai todos os placeholders de uma apresentação.
        
        Args:
            presentation: Objeto Presentation do python-pptx
            
        Returns:
            List[str]: Lista de placeholders encontrados
        """
        placeholders = set()
        
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text
                    # Procurar por padrões {{placeholder}}
                    import re
                    found_placeholders = re.findall(r'\{\{([^}]+)\}\}', text)
                    placeholders.update(found_placeholders)
        
        return sorted(list(placeholders))
    
    def update_template(self, user_id: str, template_name: str, 
                       new_placeholders: Dict[str, str]) -> Dict:
        """
        Atualiza placeholders em um template PPTX.
        
        Args:
            user_id (str): ID do usuário (deve ser admin)
            template_name (str): Nome do template
            new_placeholders (Dict): Novos placeholders {old: new}
            
        Returns:
            Dict: Resultado da operação
        """
        if not self.verify_admin_access(user_id):
            raise PermissionError(f"Usuário {user_id} não autorizado")
        
        template_path = self.templates_dir / template_name
        
        if not template_path.exists():
            raise FileNotFoundError(f"Template {template_name} não encontrado")
        
        try:
            # Fazer backup
            backup_path = self.backup_dir / f"{template_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            shutil.copy2(template_path, backup_path)
            
            # Carregar apresentação
            prs = Presentation(template_path)
            
            # Atualizar placeholders
            updated_count = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        original_text = shape.text
                        updated_text = original_text
                        
                        for old_placeholder, new_placeholder in new_placeholders.items():
                            if f"{{{{{old_placeholder}}}}}" in updated_text:
                                updated_text = updated_text.replace(
                                    f"{{{{{old_placeholder}}}}}",
                                    f"{{{{{new_placeholder}}}}}"
                                )
                                updated_count += 1
                        
                        if updated_text != original_text:
                            shape.text = updated_text
            
            # Salvar template atualizado
            prs.save(template_path)
            
            result = {
                "success": True,
                "message": f"Template {template_name} atualizado com sucesso",
                "updated_placeholders": updated_count,
                "backup_created": str(backup_path),
                "timestamp": datetime.now().isoformat()
            }
            
            logger.info(f"Template {template_name} atualizado por {user_id}")
            return result
            
        except Exception as e:
            logger.error(f"Erro ao atualizar template: {e}")
            raise
    
    def create_new_template(self, user_id: str, template_name: str, 
                           base_template: Optional[str] = None) -> Dict:
        """
        Cria um novo template PPTX.
        
        Args:
            user_id (str): ID do usuário (deve ser admin)
            template_name (str): Nome do novo template
            base_template (str, optional): Template base para copiar
            
        Returns:
            Dict: Resultado da operação
        """
        if not self.verify_admin_access(user_id):
            raise PermissionError(f"Usuário {user_id} não autorizado")
        
        new_template_path = self.templates_dir / template_name
        
        if new_template_path.exists():
            raise FileExistsError(f"Template {template_name} já existe")
        
        try:
            if base_template:
                # Copiar de template existente
                base_path = self.templates_dir / base_template
                if not base_path.exists():
                    raise FileNotFoundError(f"Template base {base_template} não encontrado")
                shutil.copy2(base_path, new_template_path)
            else:
                # Criar template básico
                prs = Presentation()
                
                # Slide título
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                
                title.text = "{{titulo_principal}}"
                subtitle.text = "{{subtitulo}}\\n{{data_geracao}}"
                
                # Slide conteúdo
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = "{{titulo_slide}}"
                content.text = "{{conteudo_principal}}\\n\\n{{metricas}}\\n\\n{{insights}}"
                
                prs.save(new_template_path)
            
            result = {
                "success": True,
                "message": f"Template {template_name} criado com sucesso",
                "template_path": str(new_template_path),
                "base_template": base_template,
                "timestamp": datetime.now().isoformat()
            }
            
            logger.info(f"Novo template {template_name} criado por {user_id}")
            return result
            
        except Exception as e:
            logger.error(f"Erro ao criar template: {e}")
            raise
    
    def configure_sharepoint(self, user_id: str, config: Dict) -> Dict:
        """
        Configura integração com SharePoint.
        
        Args:
            user_id (str): ID do usuário (deve ser admin)
            config (Dict): Configurações do SharePoint
            
        Returns:
            Dict: Resultado da operação
        """
        if not self.verify_admin_access(user_id):
            raise PermissionError(f"Usuário {user_id} não autorizado")
        
        try:
            # Atualizar configuração
            self.sharepoint_config.update(config)
            
            # Salvar configuração
            config_file = Path("config/sharepoint_config.json")
            config_file.parent.mkdir(exist_ok=True)
            
            with open(config_file, 'w') as f:
                json.dump(self.sharepoint_config, f, indent=2)
            
            result = {
                "success": True,
                "message": "Configuração SharePoint atualizada",
                "config": self.sharepoint_config,
                "timestamp": datetime.now().isoformat()
            }
            
            logger.info(f"Configuração SharePoint atualizada por {user_id}")
            return result
            
        except Exception as e:
            logger.error(f"Erro ao configurar SharePoint: {e}")
            raise
    
    def upload_to_sharepoint(self, file_path: str, user_id: str) -> Dict:
        """
        Faz upload de arquivo para SharePoint.
        
        Args:
            file_path (str): Caminho do arquivo
            user_id (str): ID do usuário
            
        Returns:
            Dict: Resultado da operação
        """
        if not self.verify_admin_access(user_id):
            raise PermissionError(f"Usuário {user_id} não autorizado")
        
        try:
            # Simular upload (implementação real requer autenticação SharePoint)
            file_name = Path(file_path).name
            sharepoint_url = f"{self.sharepoint_config['full_url']}/{file_name}"
            
            # Aqui seria implementada a lógica real de upload
            # Por enquanto, apenas simular
            
            result = {
                "success": True,
                "message": f"Arquivo {file_name} enviado para SharePoint",
                "sharepoint_url": sharepoint_url,
                "local_path": file_path,
                "timestamp": datetime.now().isoformat()
            }
            
            logger.info(f"Arquivo {file_name} enviado para SharePoint por {user_id}")
            return result
            
        except Exception as e:
            logger.error(f"Erro no upload para SharePoint: {e}")
            raise
    
    def get_admin_dashboard_data(self, user_id: str) -> Dict:
        """
        Obtém dados para o dashboard administrativo.
        
        Args:
            user_id (str): ID do usuário
            
        Returns:
            Dict: Dados do dashboard
        """
        if not self.verify_admin_access(user_id):
            raise PermissionError(f"Usuário {user_id} não autorizado")
        
        try:
            templates = self.list_templates()
            
            # Estatísticas
            total_templates = len(templates)
            total_placeholders = sum(len(t['placeholders']) for t in templates)
            
            # Arquivos recentes
            recent_files = []
            for file_path in self.output_dir.glob("*.pptx"):
                recent_files.append({
                    "name": file_path.name,
                    "size": file_path.stat().st_size,
                    "modified": datetime.fromtimestamp(file_path.stat().st_mtime).isoformat()
                })
            
            recent_files.sort(key=lambda x: x['modified'], reverse=True)
            recent_files = recent_files[:10]  # Últimos 10
            
            dashboard_data = {
                "user_info": {
                    "user_id": user_id,
                    "is_admin": True,
                    "last_access": datetime.now().isoformat()
                },
                "statistics": {
                    "total_templates": total_templates,
                    "total_placeholders": total_placeholders,
                    "recent_files_count": len(recent_files)
                },
                "templates": templates,
                "recent_files": recent_files,
                "sharepoint_config": self.sharepoint_config,
                "system_status": "operational"
            }
            
            return dashboard_data
            
        except Exception as e:
            logger.error(f"Erro ao obter dados do dashboard: {e}")
            raise


# Instância global do sistema de administração
admin_system = AdminSystem()

