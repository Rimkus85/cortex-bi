"""
Script para gerar template PPTX com placeholders para demonstração.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Criar nova apresentação
prs = Presentation()

# Slide 1: Capa
slide_layout = prs.slide_layouts[0]  # Layout de título
slide = prs.slides.add_slide(slide_layout)

# Adicionar título e subtítulo
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "{{titulo}}"
subtitle.text = "{{subtitulo}}"

# Slide 2: Comparação de Períodos
slide_layout = prs.slide_layouts[1]  # Layout de título e conteúdo
slide = prs.slides.add_slide(slide_layout)

# Adicionar título
title = slide.shapes.title
title.text = "Comparação de Períodos: {{periodo1_nome}} vs {{periodo2_nome}}"

# Adicionar tabela
left = Inches(2)
top = Inches(2)
width = Inches(6)
height = Inches(1.5)

rows, cols = 3, 3
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Cabeçalhos
table.cell(0, 0).text = "Métrica"
table.cell(0, 1).text = "{{periodo1_nome}}"
table.cell(0, 2).text = "{{periodo2_nome}}"

# Linha 1
table.cell(1, 0).text = "Total"
table.cell(1, 1).text = "{{total_periodo1}}"
table.cell(1, 2).text = "{{total_periodo2}}"

# Linha 2
table.cell(2, 0).text = "Média"
table.cell(2, 1).text = "{{media_periodo1}}"
table.cell(2, 2).text = "{{media_periodo2}}"

# Formatar cabeçalhos
for i in range(cols):
    cell = table.cell(0, i)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(14)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 3: Crescimento
slide_layout = prs.slide_layouts[1]  # Layout de título e conteúdo
slide = prs.slides.add_slide(slide_layout)

# Adicionar título
title = slide.shapes.title
title.text = "Análise de Crescimento"

# Adicionar texto com placeholders
content = slide.placeholders[1]
tf = content.text_frame

p = tf.add_paragraph()
p.text = "Crescimento Total: {{crescimento_total}}%"
p.font.size = Pt(18)

p = tf.add_paragraph()
p.text = "Crescimento Médio: {{crescimento_media}}%"
p.font.size = Pt(18)

p = tf.add_paragraph()
p.text = "Data da Análise: {{data_analise}}"
p.font.size = Pt(14)
p.font.italic = True

# Slide 4: Segmentação
slide_layout = prs.slide_layouts[1]  # Layout de título e conteúdo
slide = prs.slides.add_slide(slide_layout)

# Adicionar título
title = slide.shapes.title
title.text = "Análise de Segmentação"

# Adicionar texto com placeholders
content = slide.placeholders[1]
tf = content.text_frame

p = tf.add_paragraph()
p.text = "Total de Segmentos: {{total_segmentos}}"
p.font.size = Pt(16)

p = tf.add_paragraph()
p.text = "Total de Registros: {{total_registros}}"
p.font.size = Pt(16)

p = tf.add_paragraph()
p.text = "Maior Segmento: {{maior_segmento_nome}} ({{maior_segmento_registros}} registros, {{maior_segmento_percentual}}%)"
p.font.size = Pt(16)

# Slide 5: Motivos de Contato
slide_layout = prs.slide_layouts[1]  # Layout de título e conteúdo
slide = prs.slides.add_slide(slide_layout)

# Adicionar título
title = slide.shapes.title
title.text = "Análise de Motivos de Contato"

# Adicionar texto com placeholders
content = slide.placeholders[1]
tf = content.text_frame

p = tf.add_paragraph()
p.text = "Total de Registros: {{total_registros}}"
p.font.size = Pt(16)

p = tf.add_paragraph()
p.text = "Motivos Únicos: {{motivos_unicos}}"
p.font.size = Pt(16)

p = tf.add_paragraph()
p.text = "Motivo Mais Comum: {{motivo_mais_comum}} ({{motivo_mais_comum_count}} ocorrências, {{motivo_mais_comum_percentual}}%)"
p.font.size = Pt(16)

# Salvar apresentação
output_file = 'template_relatorio.pptx'
prs.save(output_file)

print(f"Template PPTX gerado com sucesso: {output_file}")
print("Placeholders disponíveis:")
print("- {{titulo}}, {{subtitulo}}")
print("- {{periodo1_nome}}, {{periodo2_nome}}")
print("- {{total_periodo1}}, {{total_periodo2}}")
print("- {{media_periodo1}}, {{media_periodo2}}")
print("- {{crescimento_total}}, {{crescimento_media}}")
print("- {{data_analise}}")
print("- {{total_segmentos}}, {{total_registros}}")
print("- {{maior_segmento_nome}}, {{maior_segmento_registros}}, {{maior_segmento_percentual}}")
print("- {{motivos_unicos}}, {{motivo_mais_comum}}, {{motivo_mais_comum_count}}, {{motivo_mais_comum_percentual}}")

